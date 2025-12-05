"""Service layer helpers for job statistics and automation."""

import os
import io

from django.urls import reverse
from django.utils import timezone

from accounts.models import User
from common.utils import format_currency, to_decimal
from .choices import ContentStatus, ContentSectionType
from .models import Job


def normalize_amount(value):
    return to_decimal(value)


def calculate_job_stats(queryset):
    """Return stats for cards."""

    jobs = list(queryset)
    total_jobs = len(jobs)
    pending_jobs = len([job for job in jobs if not job.is_superadmin_approved])
    total_amount = sum(normalize_amount(job.amount_inr) for job in jobs)
    return total_jobs, pending_jobs, total_amount


def get_job_cards_for_user(user):
    """Return card metadata for dashboards/welcome screens."""

    if not user.is_authenticated:
        return []

    if user.role == User.Role.MARKETING:
        jobs = Job.objects.filter(created_by=user, is_deleted__in=[False])
        total_jobs, pending_jobs, total_amount = calculate_job_stats(jobs)
        return [
            {
                "title": "Total Jobs",
                "value": total_jobs,
                "url": f"{reverse('marketing:all_projects')}?filter=all",
            },
            {
                "title": "Total Pending Jobs",
                "value": pending_jobs,
                "url": f"{reverse('marketing:all_projects')}?filter=pending",
            },
            {
                "title": "Total Amount",
                "value": format_currency(total_amount),
                "url": f"{reverse('marketing:all_projects')}?filter=amount",
            },
        ]
    if user.role == User.Role.GLOBAL:
        jobs = Job.objects.filter(created_by=user, is_deleted__in=[False])
        total_jobs, pending_jobs, total_amount = calculate_job_stats(jobs)
        return [
            {
                "title": "Total Jobs",
                "value": total_jobs,
                "url": f"{reverse('marketing:all_projects')}?filter=all",
            },
            {
                "title": "Total Pending Jobs",
                "value": pending_jobs,
                "url": f"{reverse('marketing:all_projects')}?filter=pending",
            },
            {
                "title": "Total Amount",
                "value": format_currency(total_amount),
                "url": f"{reverse('marketing:all_projects')}?filter=amount",
            },
        ]

    jobs = Job.objects.filter(is_deleted__in=[False])
    total_jobs, pending_jobs, total_amount = calculate_job_stats(jobs)
    return [
        {
            "title": "Total Jobs",
            "value": total_jobs,
            "url": reverse("superadmin:all_jobs"),
        },
        {
            "title": "Total Pending Jobs",
            "value": pending_jobs,
            "url": reverse("superadmin:new_jobs"),
        },
        {
            "title": "Total Amount",
            "value": format_currency(total_amount),
            "url": reverse("superadmin:all_jobs"),
        },
    ]


def sync_job_approval(job):
    """Mark job approved when all sections are approved."""

    if not job.sections.exists():
        job.is_superadmin_approved = False
        job.approved_at = None
        job.save(
            update_fields=["is_superadmin_approved", "approved_at", "updated_at"]
        )
        return

    all_approved = not job.sections.exclude(status=ContentStatus.APPROVED).exists()
    job.is_superadmin_approved = all_approved
    job.approved_at = timezone.now() if all_approved else None
    job.save(update_fields=["is_superadmin_approved", "approved_at", "updated_at"])


# --- Content generation helpers (Gemini-backed with fallback) --- #
# Detailed prompts adapted from the provided pipeline

def _generate_with_gemini(prompt, fallback):
    api_key = os.getenv("GEMINI_API_KEY")
    if not api_key:
        return "Generation failed: GEMINI_API_KEY not set."
    # basic retry with backoff for transient quota/rate errors
    import time
    import google.generativeai as genai

    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-2.0-flash")
    attempts = 0
    last_error = None
    while attempts < 3:
        try:
            resp = model.generate_content(prompt)
            text = getattr(resp, "text", None) or ""
            cleaned = _strip_markdown(text)
            return cleaned.strip() or fallback
        except Exception as exc:  # noqa: BLE001
            last_error = exc
            message = str(exc)
            # Retry only on obvious 429/rate/quota signals
            if "429" in message or "Resource exhausted" in message or "quota" in message.lower():
                time.sleep(2 ** attempts)  # 1s, 2s, 4s
                attempts += 1
                continue
            return f"Generation failed: {exc}" or fallback
    return f"Generation failed: {last_error}" or fallback


def _strip_markdown(text):
    """Remove common markdown headers/bullets so output stays plain text."""
    import re

    lines = (text or "").splitlines()
    cleaned = []
    for line in lines:
        stripped = line.lstrip()
        # Remove leading markdown symbols
        stripped = re.sub(r"^(#+\s*)", "", stripped)  # headings
        stripped = re.sub(r"^[-*•]\s+", "", stripped)  # bullets
        cleaned.append(stripped)
    return "\n".join(cleaned)


SUMMARY_PROMPT = """You are an AI assistant specialized in understanding writing tasks and producing a structured Job Summary, not the full content itself. Read the user's instructions and any extracted text from attachments (e.g., PDFs, DOCX) to identify what needs to be written, including topic, word count or length, reference style (APA, MLA, Harvard, etc.), and writing style or document type (essay, report, PPT, proposal, article, dissertation, thesis, etc.). If a detail is not explicitly given but can be reasonably inferred, infer it; if it cannot be inferred confidently, mark it as "Not specified." Always respond in this exact format (plain text only, no markdown symbols such as #, *, -, or bullets): Topic: <short topic or title>; Word Count: <number of words or If word count is not mentioned in the Job card, then by default print "1500">; Reference Style: <style or If Reference Style is not mentioned in the Job card, then by default print "Harvard">; Writing Style: <type or "Report">; Job Summary: <10-20 sentences clearly describing what needs to be written, the main themes to cover, target audience or level if known, and any important constraints such as tone or structure>. Do not add extra sections, do not explain your reasoning, and do not write the actual assignment—only provide a clear, concise, implementation-ready Job Summary that another writer or AI could directly follow."""

STRUCTURE_PROMPT = """You are an AI assistant specialized in creating academic writing structures (detailed outlines) for writing tasks. Your input is always the full output of a Job Summary agent, which includes at least: Topic, Word Count, Reference Style, Writing Style, and Job Summary (and may also include extra instructions). Your job is to design a clear, logically ordered, academically appropriate structure with word counts for each section and subsection, so that another writer or AI could directly draft the final document. Strictly follow all instructions and requirements from the Job Summary and ensure that every key theme, focus area, or constraint is reflected in the structure. Use academic writing conventions that match the Writing Style (e.g., essays with introduction/body/conclusion; reports with sections such as introduction, methodology, analysis, conclusion; dissertations/thesis with chapters such as introduction, literature review, methodology, results, discussion, conclusion; PPTs as slide-based academic sections, etc.). Handle Word Count as follows: always use only word counts and never pages, lines, slides, or any other length unit; if a specific word count is given, treat it as the target total and allocate section word counts so they sum to approximately that total (with minor acceptable variation); if a range is given, internally pick a reasonable midpoint and allocate based on that; if the word count is described in pages or similar, internally convert to an approximate word count and output only word counts; if Word Count is "Not specified," infer a reasonable total based on the Writing Style and academic context, then allocate accordingly. Respect the Reference Style by including a final "References" section with an appropriate word count whenever references are expected for that type of task. Ensure a coherent hierarchy with numbered sections and subsections; every primary section must include at least two numbered subheadings (e.g., 1.1, 1.2) with explicit word counts. Begin by stating the title (using the Topic) and the total word count, then list the sections in order. For every heading and subheading, provide two to four short lines that clearly describe what content should appear in that section; each line must be a standalone sentence with no bullets or symbols. Output plain text only (no markdown symbols such as #, *, -, or bullets). Do not write any actual content of the sections beyond those brief descriptive lines, only the structure, descriptions, and word counts. Do not explain your reasoning, do not add extra metadata fields, and do not mention any unit other than words."""

CONTENT_PROMPT = """You are an AI assistant specialized in academic content writing. Your input is the full output of a Structure-Making Agent, which includes the title, total word count, and a numbered list of sections and subsections with individual word counts, all derived from a Job Summary. Your task is to transform this structure into complete, polished content that strictly follows all instructions, rules, and constraints implied by both the structure and the underlying task (topic, writing style, level, focus areas, tone, etc.). You must: (1) preserve the given headings and their order exactly as provided; (2) write cohesive, formal, academic prose under each section/subsection that clearly addresses the intent of its heading and the overall task; (3) for every section and subsection, generate content that is AT LEAST the specified word count printed in the structure—treat those numbers as strict minimums per heading (e.g., if "2.1 Audit Methodology - 200 words" appears, write 200 or more words for section 2.1 alone, not counting other sections); (4) after drafting, perform an explicit internal check that every heading meets or exceeds its assigned word count AND that the combined word count of the entire response meets or exceeds the total word count stated in the structure; if any heading is short, continue writing for that heading until it meets its minimum before returning the response; (5) maintain consistency in voice, tense, and perspective as implied by the task; and (6) ensure logical flow between sections with appropriate transitions and internal coherence. Use the headings exactly as provided (no additional numbering or labels), and write the content in full paragraphs without bullets. Do not mention word counts anywhere in the content. Do not modify or invent new sections, do not change the title, and do not contradict any explicit requirements from the task (such as focus, scope, or audience). When writing the content, do not include any reference list, bibliography, or citations of any kind (no in-text citations, no author-year, no numbers in brackets, and no "References" section), even if the structure or task mentions a reference style; treat that aspect as handled elsewhere. Do not explain your reasoning or describe your process; output only the final written content organized under the given headings."""

REFERENCES_PROMPT = """You are an AI assistant specialized in generating academic reference lists and corresponding in-text citation formats. Your input will be: (1) the full content produced by a content-creation agent, (2) the specified reference style (e.g., APA, MLA, Chicago, Harvard, IEEE, etc.), and (3) the approximate total word count of the content. Your task is to create an original, topic-related reference list that strictly follows the given reference style and is based on the themes, concepts, and topics present in the content. All references you provide must be to real, credible, and verifiable sources published after 2021 (i.e., from 2022 onwards). For every 1000 words of content, generate approximately 7 references (rounding reasonably to the nearest whole number) and ensure that all references are directly relevant to the subject matter of the content. Present the references as a properly formatted “Reference List” ordered alphabetically (A–Z) by the first author’s surname, strictly conforming to the rules of the specified reference style. After the alphabetical reference list, provide a separate “Citation List” that contains the in-text citation format for each reference above (e.g., for Harvard and APA: Author, Year; for MLA: Author page; for IEEE: [number], etc.), covering all references already listed. In-text Citation rules: For Harvard, APA, APA7,  IEEE Referencing (If one, two, or three authors are present in the Reference, then use the Surname of Each Author first, then a comma, and then the year in a Single bracket). Like example: ‘Hermes, A. and Riedl, R., 2021, July. Dimensions of retail customer experience and its outcomes: a literature review and directions for future research. If you notice here, two authors are present, so the in-text citation will be “(Hermes and Riedl, 2021)”. If 4 or more authors are present, then use the first author's surname, then et al., then a comma, and then the year. For example: “Pappas, A., Fumagalli, E., Rouziou, M. and Bolander, W., 2023. More than machines: The role of the future retail salesperson in enhancing the customer experience. Journal of Retailing, 99(4), pp.518-531.”. If you notice here, 4 authors are present, so the in-text citation will be (Pappas et al. 2023). In IEEE, all are the same, but in Number Format like [1], [2], etc. Do not include any explanation, analysis, or extra text beyond the reference list and the citation list. Do not rewrite or summarize the original content. Your entire output must consist only of the formatted reference list followed by the citation list."""

FINALIZE_PROMPT = """You are an AI assistant specialized in finalizing academic documents by inserting in-text citations and appending an existing reference list. Your inputs are: (1) a complete piece of content with no citations or reference list, (2) a formatted reference list, (3) a citation list that specifies the correct in-text citation format for each reference, and (4) the reference style to follow (e.g., APA, MLA, Chicago, Harvard, IEEE, etc.). Your task is to cite all existing references from the citation list within the content and then append the full reference list at the end of the document, strictly following the given reference style. You must not rewrite, expand, shorten, reorder, or otherwise change any of the existing content, headings, or wording; you may only insert in-text citations at appropriate locations and add the reference list at the end. Don't cite in the Introduction, Conclusion parts, and if available, Abstract and Executive summary; in those parts, don't add in-text citations. Do not add new references, do not remove any existing references, and do not invent sources. Ensure that every reference from the provided reference list is cited at least once in the body using the corresponding in-text format from the citation list, and that all in-text citations match entries in the reference list. Maintain the original structure and formatting of the content as much as possible, only adding the necessary citation markers and the final reference list section. As output, return the full content with the in-text citations properly inserted and the complete reference list appended at the end, and do not include any explanations, notes, or extra commentary."""


def _extract_text_from_attachments(job):
    """Best-effort text extraction from job attachments."""
    parts = []
    for attachment in job.attachments.all():
        name = attachment.file.name.lower()
        # ensure pointer at start
        try:
            if hasattr(attachment.file, "open"):
                attachment.file.open("rb")
            attachment.file.seek(0)
        except Exception:
            pass
        try:
            data = attachment.file.read()
        except Exception:
            continue
        if name.endswith(".txt"):
            try:
                parts.append(data.decode("utf-8", errors="ignore"))
                continue
            except Exception:
                pass
        if name.endswith(".csv"):
            try:
                import pandas as pd

                df = pd.read_csv(io.BytesIO(data))
                parts.append(df.to_csv(index=False))
                continue
            except Exception:
                pass
        if name.endswith((".xlsx", ".xls", ".xlx")):
            try:
                import pandas as pd

                df = pd.read_excel(io.BytesIO(data))
                parts.append(df.to_csv(index=False))
                continue
            except Exception:
                pass
        if name.endswith(".pdf"):
            try:
                from PyPDF2 import PdfReader

                reader = PdfReader(io.BytesIO(data))
                t = "\n".join(page.extract_text() or "" for page in reader.pages)
                parts.append(t)
                continue
            except Exception:
                pass
        if name.endswith(".docx"):
            try:
                from docx import Document

                doc = Document(io.BytesIO(data))
                parts.append("\n".join(p.text for p in doc.paragraphs))
                continue
            except Exception:
                pass
        if name.endswith(".pptx"):
            try:
                from pptx import Presentation

                prs = Presentation(io.BytesIO(data))
                buff = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            buff.append(shape.text)
                parts.append("\n".join(buff))
                continue
            except Exception:
                pass
        # fallback binary decode
        try:
            parts.append(data.decode("utf-8", errors="ignore"))
        except Exception:
            parts.append(f"[Could not extract text from {attachment.file.name}]")
        try:
            attachment.file.close()
        except Exception:
            pass
    return "\n\n".join(p for p in parts if p).strip()


def generate_job_summary(job, regenerate=False, exceeded=False):
    if exceeded:
        return "Regeneration limit reached."
    instr = job.instruction or "No instruction provided."
    attachments = _extract_text_from_attachments(job)
    combined = instr
    if attachments:
        combined += "\n\nAttached file text:\n" + attachments
    base_prompt = SUMMARY_PROMPT + "\n\nUSER INPUT:\n" + combined
    return _generate_with_gemini(base_prompt.strip(), "Generation failed.")


def generate_structure_from_summary(source, regenerate=False, exceeded=False):
    """
    Generate a structure either from a Job instance or a raw summary string.

    Global flows reuse this helper with user-provided summaries, so we
    gracefully accept either datatype to ensure the stricter STRUCTURE_PROMPT
    is always used (preventing loose, emoji-heavy fallbacks).
    """

    if exceeded:
        return "Regeneration limit reached."

    if isinstance(source, str):
        summary = source.strip() or "Summary missing."
    else:
        section = (
            source.sections.filter(section_type=ContentSectionType.SUMMARY)
            .first()
        )
        summary = (section and section.content) or "Summary missing."

    base_prompt = STRUCTURE_PROMPT + "\n\nJOB SUMMARY:\n" + summary
    return _generate_with_gemini(base_prompt.strip(), "Structure generation failed.")


def generate_content_from_structure(source, regenerate=False, exceeded=False):
    if exceeded:
        return "Regeneration limit reached."
    if isinstance(source, str):
        structure = source.strip() or "Structure missing."
    else:
        section = (
            source.sections.filter(section_type=ContentSectionType.STRUCTURE).first()
        )
        structure = (section and section.content) or "Structure missing."
    base_prompt = CONTENT_PROMPT + "\n\nSTRUCTURE:\n" + structure
    return _generate_with_gemini(base_prompt.strip(), "Content generation failed.")


def generate_references_from_content(job, regenerate=False, exceeded=False):
    if exceeded:
        return "Regeneration limit reached."
    content = (
        job.sections.filter(section_type=ContentSectionType.CONTENT).first().content
        or "Content missing."
    )
    base_prompt = REFERENCES_PROMPT + "\n\nCONTENT:\n" + content
    return _generate_with_gemini(base_prompt.strip(), "Generation failed.")


def generate_final_document_with_citations(job, regenerate=False, exceeded=False):
    if exceeded:
        return "Regeneration limit reached."
    refs = (
        job.sections.filter(section_type=ContentSectionType.REFERENCING)
        .first()
        .content
        or "References missing."
    )
    content = (
        job.sections.filter(section_type=ContentSectionType.CONTENT)
        .first()
        .content
        or "Content missing."
    )
    combined = (
        FINALIZE_PROMPT
        + "\n\n=== CONTENT (NO CITATIONS) ===\n"
        + content
        + "\n\n=== REFERENCES ===\n"
        + refs
    )
    return _generate_with_gemini(combined.strip(), "Generation failed.")
